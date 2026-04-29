import io
import logging
import os
import traceback
from typing import Any

from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import HTMLResponse, Response
from fastapi.staticfiles import StaticFiles
from pptx import Presentation
from pptx.oxml.ns import qn
import httpx
from pydantic import BaseModel

logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")
log = logging.getLogger(__name__)

app = FastAPI(title="JeoparTy Generator API")

@app.get("/health")
def health():
    log.info("Health check hit")
    return {"status": "ok"}

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["POST", "GET"],
    allow_headers=["*"],
)

HERE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_PATH = os.path.join(HERE, "template.pptm")
POINT_VALUES = [200, 400, 600, 800, 1000]
MASTER_API_KEY = os.environ.get("GEMINI_API_KEY", "")

# Slides to remove for simple (no scoreboard) version — 0-based indices
# PowerPoint slides 2-3   → indices 1-2
# PowerPoint slides 60-61 → indices 59-60
# PowerPoint slides 65-70 → indices 64-69
SIMPLE_SLIDES_TO_REMOVE = [1, 2, 59, 60, 64, 65, 66, 67, 68, 69]


class GenerateRequest(BaseModel):
    game_data: dict[str, Any]
    theme: str = "Game"
    api_key: str = ""
    version: str = "full"


# ── Helper functions ──────────────────────────────────────────────────────────

def set_shape_text(slide, shape_name: str, text: str) -> bool:
    for shape in slide.shapes:
        if shape.name != shape_name:
            continue
        if not shape.has_text_frame:
            log.warning("Shape %s has no text frame", shape_name)
            return False
        tf = shape.text_frame
        if not tf.paragraphs:
            log.warning("Shape %s has no paragraphs", shape_name)
            return False
        para = tf.paragraphs[0]
        if not para.runs:
            run = para.add_run()
            run.text = text
        else:
            para.runs[0].text = text
            for run in para.runs[1:]:
                run.text = ""
        return True
    log.warning("Shape not found: %s", shape_name)
    return False


def remove_vba(prs: Presentation) -> None:
    """Strip the VBA project from the presentation so it saves cleanly as .pptx.
    Finds the vbaProject relationship, drops it, and removes the part reference
    from the presentation XML so PowerPoint sees a clean .pptx."""
    VBA_REL = 'http://schemas.microsoft.com/office/2006/relationships/vbaProject'
    part = prs.part
    # Find the rId for the vbaProject relationship
    rId_to_drop = None
    for rId, rel in part.rels.items():
        if rel.reltype == VBA_REL:
            rId_to_drop = rId
            break
    if rId_to_drop:
        part.drop_rel(rId_to_drop)
        log.info("VBA project removed (rId: %s)", rId_to_drop)
    else:
        log.info("No VBA project found — nothing to remove")


def delete_slides(prs: Presentation, indices: list[int]) -> None:
    """Delete slides by 0-based index. Must delete in reverse order
    to avoid index shifting as slides are removed."""
    slide_ids = prs.slides._sldIdLst
    for idx in sorted(set(indices), reverse=True):
        if idx >= len(prs.slides):
            log.warning("Slide index %d out of range (%d slides), skipping", idx, len(prs.slides))
            continue
        rId = slide_ids[idx].get(qn('r:id'))
        if rId:
            prs.part.drop_rel(rId)
        del slide_ids[idx]
        log.info("Deleted slide index %d", idx)


# ── API endpoints ─────────────────────────────────────────────────────────────

@app.get("/has-master-key")
async def has_master_key():
    """Let the frontend know if a master key is available."""
    return {"available": bool(MASTER_API_KEY)}


class GeminiRequest(BaseModel):
    prompt: str


@app.post("/gemini")
async def call_gemini(req: GeminiRequest):
    """Call Gemini using the master key stored in environment."""
    if not MASTER_API_KEY:
        raise HTTPException(400, "No master API key configured on server")

    url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent?key={MASTER_API_KEY}"
    payload = {
        "contents": [{"parts": [{"text": req.prompt}]}],
        "generationConfig": {"temperature": 0.7, "maxOutputTokens": 8192}
    }

    async with httpx.AsyncClient(timeout=120) as client:
        res = await client.post(url, json=payload)
    if res.status_code != 200:
        raise HTTPException(502, f"Gemini API error: {res.text[:200]}")
    data = res.json()
    text = data.get("candidates", [{}])[0].get("content", {}).get("parts", [{}])[0].get("text", "")
    if not text:
        raise HTTPException(502, "Empty response from Gemini")
    log.info("Gemini via master key returned %d chars", len(text))
    return {"text": text}


@app.post("/generate")
async def generate(req: GenerateRequest):
    data = req.game_data
    version = req.version.strip().lower()

    for key in ["categories", "clues", "answers",
                "finalJeopardyClue", "finalJeopardyAnswer", "finalJeopardyTopic"]:
        if key not in data:
            raise HTTPException(400, f"Missing field: {key}")

    categories = data["categories"]
    clues      = data["clues"]
    answers    = data["answers"]

    log.info("Categories: %s", categories)
    log.info("Version: %s", version)

    if len(categories) != 5:
        raise HTTPException(400, f"Expected 5 categories, got {len(categories)}")
    if len(clues) != 5 or len(answers) != 5:
        raise HTTPException(400, "Expected 5 difficulty rows")

    if not os.path.exists(TEMPLATE_PATH):
        raise HTTPException(500, f"Template not found: {TEMPLATE_PATH}")

    try:
        prs = Presentation(TEMPLATE_PATH)
        log.info("Template loaded. Slides: %d", len(prs.slides))

        # ── 1. Main panel category titles (Slide 1, index 0) ──
        panel_slide = prs.slides[0]
        for i, cat in enumerate(categories, start=1):
            ok = set_shape_text(panel_slide, f"Title_Cat{i}", cat)
            log.info("Title_Cat%d -> '%s' [%s]", i, cat, "OK" if ok else "MISS")

        # ── 2. Hidden data slides — found by name, not index ──
        def find_slide_by_name(name):
            for s in prs.slides:
                if s.name == name:
                    return s
            raise HTTPException(500, f"Data slide not found: '{name}'")

        q_data_slide  = find_slide_by_name("DataSlide_Questions")
        a_data_slide  = find_slide_by_name("DataSlide_Answers")
        fj_data_slide = find_slide_by_name("DataSlide_FinalJeopardy")

        for i, cat in enumerate(categories, start=1):
            set_shape_text(q_data_slide, f"Data_Cat{i}", cat)

        # ── 3. Game question & answer slides ──
        for diff_idx, points in enumerate(POINT_VALUES):
            for cat_idx in range(5):
                cat_num = cat_idx + 1
                clue   = clues[diff_idx][cat_idx]   if diff_idx < len(clues)   and cat_idx < len(clues[diff_idx])   else ""
                answer = answers[diff_idx][cat_idx] if diff_idx < len(answers) and cat_idx < len(answers[diff_idx]) else ""

                q_shape_name = f"Q_Cat{cat_num}_{points}"
                a_shape_name = f"A_Cat{cat_num}_{points}"

                q_found = a_found = False
                for slide in prs.slides:
                    if not q_found:
                        q_found = set_shape_text(slide, q_shape_name, str(clue))
                    if not a_found:
                        a_found = set_shape_text(slide, a_shape_name, str(answer))
                    if q_found and a_found:
                        break

                set_shape_text(q_data_slide, f"Data_Q_Cat{cat_num}_{points}", str(clue))
                set_shape_text(a_data_slide, f"Data_A_Cat{cat_num}_{points}", str(answer))

                if not q_found or not a_found:
                    log.warning("MISS: %s=%s %s=%s", q_shape_name, q_found, a_shape_name, a_found)
                else:
                    log.info("✓ Cat%d %d: '%s' / '%s'", cat_num, points, str(clue)[:40], str(answer))

        # ── 4. Final Jeopardy ──
        fj_topic  = str(data["finalJeopardyTopic"])
        fj_clue   = str(data["finalJeopardyClue"])
        fj_answer = str(data["finalJeopardyAnswer"])

        for slide in prs.slides:
            set_shape_text(slide, "FinalJeopardyTopic", fj_topic)
            set_shape_text(slide, "FinalJeopardyClue",  fj_clue)
            set_shape_text(slide, "FinalJeopardyAnswer", fj_answer)

        set_shape_text(fj_data_slide, "Data_FJ_Topic",  fj_topic)
        set_shape_text(fj_data_slide, "Data_FJ_Clue",   fj_clue)
        set_shape_text(fj_data_slide, "Data_FJ_Answer",  fj_answer)

        log.info("Final Jeopardy: %s / %s", fj_topic, fj_answer)

        # ── 5. Build file ──
        slug = "".join(c for c in str(req.theme)[:20] if c.isalnum() or c in " -_").strip() or "Game"

        if version == "simple":
            log.info("Simple version — stripping VBA and removing scoreboard slides")
            remove_vba(prs)
            delete_slides(prs, SIMPLE_SLIDES_TO_REMOVE)
            log.info("Slides remaining: %d", len(prs.slides))
            media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            filename = f"AI Jeopardy - {slug}.pptx"
        else:
            media_type = "application/vnd.ms-powerpoint.presentation.macroEnabled.12"
            filename = f"AI Jeopardy - {slug}.pptm"

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        file_bytes = buf.read()
        log.info("File built: %d bytes, version: %s", len(file_bytes), version)

        return Response(
            content=file_bytes,
            media_type=media_type,
            headers={"Content-Disposition": f'attachment; filename="{filename}"'},
        )

    except HTTPException:
        raise
    except Exception as e:
        log.error("Generation failed: %s", traceback.format_exc())
        raise HTTPException(500, f"Failed to generate file: {str(e)}")


# ── Static files and root ─────────────────────────────────────────────────────

static_dir = os.path.join(HERE, "static")
if os.path.exists(static_dir):
    app.mount("/static", StaticFiles(directory=static_dir), name="static")


@app.get("/", response_class=HTMLResponse)
async def root():
    html_path = os.path.join(HERE, "index.html")
    if not os.path.exists(html_path):
        raise HTTPException(404, "index.html not found")
    with open(html_path) as f:
        return f.read()
